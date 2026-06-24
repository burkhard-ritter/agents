#!/usr/bin/env python3
"""
Create a new PowerPoint presentation with title and content slides.

Usage:
    uv run --with python-pptx create_presentation.py <output_file> <title> [subtitle]

Example:
    uv run --with python-pptx create_presentation.py report.pptx "Q4 Results" "2025 Performance"
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt


def create_presentation(output_file, title, subtitle=""):
    """Create a new presentation with a title slide."""
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Save presentation
    prs.save(output_file)
    print(f"✓ Created presentation: {output_file}")
    print(f"  Title: {title}")
    if subtitle:
        print(f"  Subtitle: {subtitle}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: create_presentation.py <output_file> <title> [subtitle]")
        print("\nExample:")
        print('  uv run --with python-pptx create_presentation.py report.pptx "Q4 Results" "2025"')
        sys.exit(1)

    output_file = sys.argv[1]
    title = sys.argv[2]
    subtitle = sys.argv[3] if len(sys.argv) > 3 else ""

    create_presentation(output_file, title, subtitle)
