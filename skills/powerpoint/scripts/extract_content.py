#!/usr/bin/env python3
"""
Extract text content and metadata from a PowerPoint presentation.

Usage:
    uv run --with python-pptx extract_content.py <input_file> [--format json|text]

Example:
    uv run --with python-pptx extract_content.py presentation.pptx
    uv run --with python-pptx extract_content.py presentation.pptx --format json
"""

import sys
import json
from pptx import Presentation


def extract_text(prs):
    """Extract all text from a presentation."""
    content = []

    for i, slide in enumerate(prs.slides):
        slide_content = {
            'slide_number': i + 1,
            'shapes': []
        }

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                shape_data = {
                    'text': shape.text,
                    'type': shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type)
                }
                slide_content['shapes'].append(shape_data)

        if slide_content['shapes']:
            content.append(slide_content)

    return content


def print_text_format(content):
    """Print content in human-readable text format."""
    for slide in content:
        print(f"\n{'='*60}")
        print(f"Slide {slide['slide_number']}")
        print(f"{'='*60}")

        for shape in slide['shapes']:
            print(f"\n[{shape['type']}]")
            print(shape['text'])

    print(f"\n{'='*60}")
    print(f"Total slides: {len(content)}")


def print_json_format(content):
    """Print content in JSON format."""
    print(json.dumps(content, indent=2, ensure_ascii=False))


def extract_metadata(prs):
    """Extract presentation metadata."""
    core_props = prs.core_properties

    metadata = {
        'title': core_props.title or "",
        'author': core_props.author or "",
        'subject': core_props.subject or "",
        'created': core_props.created.isoformat() if core_props.created else "",
        'modified': core_props.modified.isoformat() if core_props.modified else "",
        'slide_count': len(prs.slides),
        'slide_width': prs.slide_width,
        'slide_height': prs.slide_height,
    }

    return metadata


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: extract_content.py <input_file> [--format json|text]")
        print("\nExample:")
        print("  uv run --with python-pptx extract_content.py presentation.pptx")
        sys.exit(1)

    input_file = sys.argv[1]
    output_format = 'text'

    if len(sys.argv) > 2 and sys.argv[2] == '--format':
        if len(sys.argv) > 3:
            output_format = sys.argv[3]

    try:
        prs = Presentation(input_file)

        # Extract metadata
        metadata = extract_metadata(prs)

        # Extract content
        content = extract_text(prs)

        if output_format == 'json':
            result = {
                'metadata': metadata,
                'content': content
            }
            print_json_format(result)
        else:
            print(f"\nPresentation: {input_file}")
            print(f"Title: {metadata['title']}")
            print(f"Author: {metadata['author']}")
            print(f"Slides: {metadata['slide_count']}")
            print_text_format(content)

    except FileNotFoundError:
        print(f"Error: File '{input_file}' not found")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
