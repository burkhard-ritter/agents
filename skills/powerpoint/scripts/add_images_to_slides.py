#!/usr/bin/env python3
"""
Add background or side images to PowerPoint slides.

Usage:
    uv run --with python-pptx --with requests add_images_to_slides.py <input_file> <output_file> <config.json>

Example:
    uv run --with python-pptx --with requests add_images_to_slides.py input.pptx output.pptx image_config.json
"""

import sys
import json
import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def download_image(url, filename):
    """Download an image from a URL."""
    try:
        response = requests.get(url, stream=True, timeout=30)
        response.raise_for_status()

        with open(filename, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)

        print(f"✓ Downloaded: {filename}")
        return True
    except Exception as e:
        print(f"✗ Failed to download {url}: {e}")
        return False


def add_background_image(slide, prs, image_path, opacity=0.5):
    """Add a full background image to a slide."""
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Add image as background (full slide)
    pic = slide.shapes.add_picture(
        image_path,
        0, 0,  # left, top
        width=slide_width,
        height=slide_height
    )

    # Send to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)

    # Add semi-transparent overlay for readability
    overlay = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0,
        slide_width,
        slide_height
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
    overlay.fill.transparency = opacity
    overlay.line.fill.background()

    # Send overlay behind text but over image
    slide.shapes._spTree.remove(overlay._element)
    slide.shapes._spTree.insert(3, overlay._element)


def add_side_image(slide, prs, image_path, side='right', width_percent=35):
    """Add an image to the side of a slide."""
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Calculate image dimensions
    img_width = int(slide_width * (width_percent / 100))

    if side == 'right':
        left = slide_width - img_width
    else:  # left
        left = 0

    # Add image
    pic = slide.shapes.add_picture(
        image_path,
        left, 0,
        width=img_width,
        height=slide_height
    )

    # Send to back
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def process_presentation(input_file, output_file, config):
    """Process presentation and add images based on configuration."""
    prs = Presentation(input_file)

    images_dir = "downloaded_images"
    os.makedirs(images_dir, exist_ok=True)

    for slide_config in config.get('slides', []):
        slide_num = slide_config['slide_number']

        if slide_num > len(prs.slides):
            print(f"⚠ Slide {slide_num} doesn't exist, skipping")
            continue

        slide = prs.slides[slide_num - 1]  # 0-indexed

        image_url = slide_config.get('image_url')
        image_type = slide_config.get('type', 'side')  # 'side' or 'background'

        if not image_url:
            print(f"⚠ No image URL for slide {slide_num}, skipping")
            continue

        # Check if image_url is a local file or a URL
        if os.path.exists(image_url):
            # It's a local file
            image_filename = image_url
            print(f"Using local image: {image_filename}")
        elif image_url.startswith(('http://', 'https://')):
            # It's a URL, download it
            image_filename = os.path.join(images_dir, f"slide_{slide_num}.jpg")

            if not os.path.exists(image_filename):
                print(f"Downloading image for slide {slide_num}...")
                if not download_image(image_url, image_filename):
                    continue
        else:
            print(f"⚠ Invalid image path/URL for slide {slide_num}: {image_url}")
            continue

        # Add image to slide
        print(f"Adding {image_type} image to slide {slide_num}...")

        if image_type == 'background':
            opacity = slide_config.get('opacity', 0.5)
            add_background_image(slide, prs, image_filename, opacity)
        else:  # side
            side = slide_config.get('side', 'right')
            width = slide_config.get('width_percent', 35)
            add_side_image(slide, prs, image_filename, side, width)

        print(f"✓ Processed slide {slide_num}")

    # Save modified presentation
    prs.save(output_file)
    print(f"\n✓ Saved modified presentation: {output_file}")


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: add_images_to_slides.py <input_file> <output_file> <config.json>")
        print("\nExample config.json:")
        print("""{
  "slides": [
    {
      "slide_number": 3,
      "type": "side",
      "side": "right",
      "width_percent": 35,
      "image_url": "https://images.unsplash.com/photo-..."
    },
    {
      "slide_number": 4,
      "type": "background",
      "opacity": 0.6,
      "image_url": "https://images.unsplash.com/photo-..."
    }
  ]
}""")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2]
    config_file = sys.argv[3]

    try:
        with open(config_file, 'r') as f:
            config = json.load(f)

        process_presentation(input_file, output_file, config)

    except FileNotFoundError as e:
        print(f"Error: File not found - {e}")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in config file - {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)
