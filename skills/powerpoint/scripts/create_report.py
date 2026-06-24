#!/usr/bin/env python3
"""
Create a report-style presentation from structured data.

Usage:
    uv run --with python-pptx create_report.py <output_file> <data_file.json>

The JSON file should have this structure:
{
    "title": "Report Title",
    "subtitle": "Report Subtitle",
    "sections": [
        {
            "title": "Section 1",
            "slides": [
                {
                    "title": "Slide Title",
                    "content": ["Bullet 1", "Bullet 2"]
                }
            ]
        }
    ]
}

Example:
    uv run --with python-pptx create_report.py quarterly_report.pptx data.json
"""

import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt


def create_report(output_file, data):
    """Create a report presentation from structured data."""
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = data.get('title', 'Report')
    if 'subtitle' in data:
        slide.placeholders[1].text = data['subtitle']

    # Add sections and slides
    for section in data.get('sections', []):
        # Add section header
        section_slide_layout = prs.slide_layouts[2]
        slide = prs.slides.add_slide(section_slide_layout)
        slide.shapes.title.text = section.get('title', 'Section')

        # Add content slides
        for slide_data in section.get('slides', []):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            slide.shapes.title.text = slide_data.get('title', '')

            content = slide_data.get('content', [])
            if content:
                text_frame = slide.placeholders[1].text_frame
                text_frame.clear()

                for i, bullet in enumerate(content):
                    if i == 0:
                        text_frame.text = bullet
                    else:
                        p = text_frame.add_paragraph()
                        p.text = bullet
                        p.level = 0

    prs.save(output_file)
    return prs


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: create_report.py <output_file> <data_file.json>")
        print("\nExample:")
        print("  uv run --with python-pptx create_report.py quarterly_report.pptx data.json")
        sys.exit(1)

    output_file = sys.argv[1]
    data_file = sys.argv[2]

    try:
        with open(data_file, 'r') as f:
            data = json.load(f)

        prs = create_report(output_file, data)

        print(f"✓ Created report presentation: {output_file}")
        print(f"  Slides: {len(prs.slides)}")

    except FileNotFoundError:
        print(f"Error: File '{data_file}' not found")
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in {data_file}: {e}")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
