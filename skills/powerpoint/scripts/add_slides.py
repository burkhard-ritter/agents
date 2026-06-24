#!/usr/bin/env python3
"""
Add slides to an existing PowerPoint presentation.

Usage:
    uv run --with python-pptx add_slides.py <input_file> <slide_type> <title> [content]

Slide types: title, bullet, blank, section

Example:
    uv run --with python-pptx add_slides.py report.pptx bullet "Key Points" "First point"
    uv run --with python-pptx add_slides.py report.pptx section "Q2 Results"
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

    return slide


def add_bullet_slide(prs, title, content=""):
    """Add a bullet point slide."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    slide.shapes.title.text = title

    if content:
        text_frame = slide.placeholders[1].text_frame
        text_frame.text = content

    return slide


def add_section_slide(prs, title):
    """Add a section header slide."""
    section_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(section_slide_layout)

    slide.shapes.title.text = title

    return slide


def add_blank_slide(prs):
    """Add a blank slide."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    return slide


if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: add_slides.py <input_file> <slide_type> <title> [content]")
        print("\nSlide types: title, bullet, blank, section")
        print("\nExamples:")
        print('  uv run --with python-pptx add_slides.py report.pptx bullet "Key Points" "First point"')
        print('  uv run --with python-pptx add_slides.py report.pptx section "Q2 Results"')
        sys.exit(1)

    input_file = sys.argv[1]
    slide_type = sys.argv[2].lower()
    title = sys.argv[3] if len(sys.argv) > 3 else ""
    content = sys.argv[4] if len(sys.argv) > 4 else ""

    try:
        prs = Presentation(input_file)

        if slide_type == 'title':
            add_title_slide(prs, title, content)
        elif slide_type == 'bullet':
            add_bullet_slide(prs, title, content)
        elif slide_type == 'section':
            add_section_slide(prs, title)
        elif slide_type == 'blank':
            add_blank_slide(prs)
        else:
            print(f"Error: Unknown slide type '{slide_type}'")
            print("Valid types: title, bullet, blank, section")
            sys.exit(1)

        prs.save(input_file)
        print(f"✓ Added {slide_type} slide to {input_file}")
        if title:
            print(f"  Title: {title}")

    except FileNotFoundError:
        print(f"Error: File '{input_file}' not found")
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
